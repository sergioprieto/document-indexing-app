import boto3
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from llama_index.core import Document
from llama_index.vector_stores.postgres import PGVectorStore
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.core.node_parser.text import SentenceSplitter
import urllib.parse
import os
from typing import List, Dict, Optional
import re
import asyncio
import json
import time
import openai

# Read environment variables
OPENAI_API_KEY = os.environ['OPENAI_API_KEY']
PG_HOST = os.environ['PG_HOST']
PG_PORT = os.environ['PG_PORT']
PG_DATABASE = os.environ['PG_DATABASE']
PG_USER = os.environ['PG_USER']
PG_PASSWORD = os.environ['PG_PASSWORD']
ASYNC_CONNECTION_STRING = f"postgresql+asyncpg://{PG_USER}:{PG_PASSWORD}@{PG_HOST}:{PG_PORT}/{PG_DATABASE}"
SCHEMA_NAME = "public"
DEFAULT_CHUNK_SIZE = 850
SENTENCE_CHUNK_OVERLAP = 50

# read openai key
openai.api_key = OPENAI_API_KEY

s3 = boto3.client('s3')
embedding_model = OpenAIEmbedding(model="text-embedding-ada-002")

# HELPER FUNCTIONS
def process_file(file_path: str) -> str:
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_extension == ".pdf":
            with open(file_path, "rb") as file:
                pdf_reader = PdfReader(file)
                content = ""
                for page in pdf_reader.pages:
                    content += page.extract_text()
        elif file_extension == ".docx":
            docx_document = DocxDocument(file_path)
            content = "\n".join([paragraph.text for paragraph in docx_document.paragraphs])
        elif file_extension == ".pptx":
            pptx_document = Presentation(file_path)
            content = ""
            for slide in pptx_document.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        elif file_extension == ".xlsx":
            workbook = load_workbook(file_path)
            content = ""
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        if cell.value is not None:
                            row_data.append(str(cell.value))
                    if row_data:
                        content += " ".join(row_data) + "\n"
        elif file_extension == ".txt":
            with open(file_path, "r") as file:
                content = file.read()
        else:
            print(f"Unsupported file format: {file_extension}")
            return ""
    except Exception as e:
        print(f"Error processing file: {str(e)}")
        return ""
    
    return content

def extract_links(text: str) -> List[str]:
    top_tlds = {
        "com", "net", "org", "jp", "de", "uk", "fr", "br", "it", "ru", "es", "me", "gov", "pl", 
        "ca", "au", "cn", "co", "in", "nl", "edu", "info", "eu", "ch", "id", "at", "kr", "cz", 
        "mx", "be", "tv", "se", "tr", "tw", "al", "ua", "ir", "vn", "cl", "sk", "ly", "cc", 
        "to", "no", "fi", "us", "pt", "dk", "ar", "hu", "tk", "gr", "il", "news", "ro", "my", 
        "biz", "ie", "za", "nz", "sg", "ee", "th", "io", "xyz", "pe", "bg", "hk", "rs", "lt", 
        "link", "ph", "club", "si", "site", "mobi", "by", "cat", "wiki", "la", "ga", "xxx", 
        "cf", "hr", "ng", "jobs", "online", "kz", "ug", "gq", "ae", "is", "lv", "pro", "fm", 
        "tips", "ms", "sa", "app"
    }

    sorted_tlds = sorted(top_tlds, key=len, reverse=True)
    tld_pattern = '|'.join(re.escape(tld) for tld in sorted_tlds)
    
    url_pattern = re.compile(rf'\b(https?://(?:(?:[\w-]+\.)+(?:{tld_pattern}))(?:/[^\s.]*)?(?<!\.))', re.IGNORECASE)

    links = re.findall(url_pattern, text)

    return links

class LoggingSentenceSplitter(SentenceSplitter):
    def split_text(self, text: str, metadata: Optional[Dict] = None) -> List[str]:
        chunks = super().split_text(text, metadata)
        for i, chunk in enumerate(chunks):
            print(f"Chunk {i} metadata: {json.dumps(metadata, indent=2)}")
            print(f"Chunk {i} length: {len(chunk)}")
        return chunks

def simplify_links(links, max_subdomains_per_domain=5):
    from urllib.parse import urlparse
    from collections import defaultdict

    simplified = set()
    domain_subdomains = defaultdict(set)

    for link in links:
        parsed_url = urlparse(link)
        domain = f"{parsed_url.scheme}://{parsed_url.netloc}"
        path = parsed_url.path

        # Add the main domain
        simplified.add(domain)

        # If there's a path, consider it as a subdomain
        if path and path != '/':
            # Only add if we haven't reached the max subdomains for this domain
            if len(domain_subdomains[domain]) < max_subdomains_per_domain:
                domain_subdomains[domain].add(link)

    # Add the selected subdomains to the simplified set
    for subdomains in domain_subdomains.values():
        simplified.update(subdomains)

    return list(simplified)

def load_documents_from_file(file_path: str) -> List[Document]:
    documents = []
    file_name = os.path.basename(file_path)
    
    content = process_file(file_path)
    if content:
        links = extract_links(content)
        
        # Simplify and reduce links
        simplified_links = simplify_links(links, max_subdomains_per_domain=5)
        
        document = Document(
            text=content,
            doc_id=file_name,
            metadata={
                "file_id": file_name,
                "source": file_name,
                "links": simplified_links
            }
        )
        documents.append(document)
        
        print(f"Document metadata: {json.dumps(document.metadata, indent=2)}")
        print(f"Number of links before simplification: {len(links)}")
        print(f"Number of links after simplification: {len(simplified_links)}")
    
    return documents

async def process_event(event, context):
    print("JOB STARTED")
    start_time = time.time()
    vector_store = None

    try:
        bucket_name = event['Records'][0]['s3']['bucket']['name']
        object_key = urllib.parse.unquote_plus(event['Records'][0]['s3']['object']['key'])
        event_name = event['Records'][0]['eventName']
        print(f"Processing file: {object_key} from S3 bucket: {bucket_name}")
        print(f"Event Name: {event_name}")

        parent_folder = os.path.dirname(object_key)
        table_name = f"{parent_folder}_vector_index".replace('/', '-')

        vector_store = PGVectorStore.from_params(
            host=PG_HOST,
            port=PG_PORT,
            database=PG_DATABASE,
            user=PG_USER,
            password=PG_PASSWORD,
            table_name=table_name,
            embed_dim=1536,
            async_connection_string=ASYNC_CONNECTION_STRING
        )

        # Check if _initialize method exists and call it appropriately
        if hasattr(vector_store, '_initialize'):
            if asyncio.iscoroutinefunction(vector_store._initialize):
                await vector_store._initialize()
            else:
                vector_store._initialize()
        else:
            print("Warning: _initialize method not found")

        if event_name.startswith("ObjectCreated") or event_name.startswith("ObjectModified"):
            local_file_path = '/tmp/' + os.path.basename(object_key)
            try:
                s3.download_file(bucket_name, object_key, local_file_path)
                print(f"File downloaded to: {local_file_path}")
                
                documents = load_documents_from_file(local_file_path)
                print(f"Loaded {len(documents)} documents")
                
                node_parser = SentenceSplitter(
                    chunk_size=DEFAULT_CHUNK_SIZE,
                    chunk_overlap=SENTENCE_CHUNK_OVERLAP,
                )
                
                all_nodes = []
                for doc in documents:
                    nodes = node_parser.get_nodes_from_documents([doc])
                    all_nodes.extend(nodes)
                
                print(f"Created {len(all_nodes)} nodes from {len(documents)} documents.")

                embedding_start = time.time()
                for node in all_nodes:
                    node.embedding = embedding_model.get_text_embedding(node.get_content())
                
                embedding_end = time.time()
                print(f"Embedded {len(all_nodes)} nodes in {embedding_end - embedding_start:.2f} seconds.")

                insert_start = time.time()
                if hasattr(vector_store, 'async_add'):
                    if asyncio.iscoroutinefunction(vector_store.async_add):
                        inserted_ids = await vector_store.async_add(all_nodes)
                    else:
                        inserted_ids = vector_store.async_add(all_nodes)
                elif hasattr(vector_store, 'add'):
                    inserted_ids = vector_store.add(all_nodes)
                else:
                    raise AttributeError("No suitable add method found on vector_store")
                
                insert_end = time.time()
                print(f"Added {len(inserted_ids)} new nodes to the vector store in {insert_end - insert_start:.2f} seconds.")

            except Exception as e:
                print(f"Error processing file: {str(e)}")
                raise
    
        elif event_name.startswith("ObjectRemoved"):
            try:
                file_name = os.path.basename(object_key)
                delete_start = time.time()
                if hasattr(vector_store, 'delete'):
                    if asyncio.iscoroutinefunction(vector_store.delete):
                        await vector_store.delete(file_name)
                    else:
                        vector_store.delete(file_name)
                    delete_end = time.time()
                    print(f"Deleted documents from vector store for file: {file_name} in {delete_end - delete_start:.2f} seconds")
                else:
                    print("Warning: delete method not found on vector_store")
            except Exception as e:
                print(f"Error deleting documents from vector store: {str(e)}")
                raise

        else:
            print(f"Unsupported event type: {event_name}")

    except Exception as e:
        print(f"Error processing event: {str(e)}")
        raise
    finally:
        if vector_store:
            if hasattr(vector_store, 'close'):
                if asyncio.iscoroutinefunction(vector_store.close):
                    await vector_store.close()
                else:
                    vector_store.close()
            else:
                print("Warning: close method not found")

    end_time = time.time()
    print(f"JOB FINISHED in {end_time - start_time:.2f} seconds")

    return {
        'statusCode': 200,
        'body': 'END OF LAMBDA FUNCTION'
    }

def handler(event, context):
    return asyncio.run(process_event(event, context))